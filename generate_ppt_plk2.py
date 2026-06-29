"""
Script generate PPT Presentasi Sidang PLK 2 - Farrel Aqeel Danendra
Rekognisi: Analisis Kebutuhan | Pemrograman API | Uji Coba & Implementasi
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Konstanta Warna ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
DARK       = RGBColor(0x1F, 0x29, 0x37)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN      = RGBColor(0x05, 0x96, 0x69)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLD_W = prs.slide_width
SLD_H = prs.slide_height

# ── Helper Functions ─────────────────────────────────────────────────────────
def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_image_placeholder(slide, left, top, width, height, label="[MASUKKAN GAMBAR DI SINI]"):
    """Add a placeholder box for images."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = GRAY
    shape.line.dash_style = 2  # dashed
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(height // 12700 // 2 - 10)
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK):
    """Add a bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(4)
        p.level = 0
    return txBox

def add_code_block(slide, left, top, width, height, code_text, title=""):
    """Add a code block with dark background."""
    # Background
    shape = add_rect(slide, left, top, width, height, RGBColor(0x1E, 0x29, 0x3B))
    # Title bar
    if title:
        title_shape = add_rect(slide, left, top, width, Inches(0.35), RGBColor(0x0F, 0x17, 0x2A))
        add_text(slide, left + Inches(0.2), top + Inches(0.05), width - Inches(0.4), Inches(0.25),
                 title, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8), bold=True)
        code_top = top + Inches(0.4)
        code_h = height - Inches(0.4)
    else:
        code_top = top + Inches(0.1)
        code_h = height - Inches(0.2)
    # Code text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), code_top, width - Inches(0.4), code_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.split('\n')[:30]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p.font.name = 'Consolas'
        p.space_before = Pt(1)
    return shape

def slide_number_footer(slide, num, total):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
             f"{num}/{total}", font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def section_header(slide, section_num, title, subtitle=""):
    """Add section divider slide."""
    add_bg(slide, NAVY)
    # Accent bar
    add_rect(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(1.5), BLUE)
    add_text(slide, Inches(1.2), Inches(2.5), Inches(2), Inches(0.5),
             section_num, font_size=14, color=BLUE, bold=True)
    add_text(slide, Inches(1.2), Inches(3.0), Inches(8), Inches(1.0),
             title, font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(1.2), Inches(4.2), Inches(8), Inches(0.6),
                 subtitle, font_size=16, color=GRAY)

TOTAL_SLIDES = 24

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
# Decorative elements
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)
add_rect(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), BLUE)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
         "LAPORAN AKHIR — MAGANG MANDIRI", font_size=14, color=BLUE, bold=True)
add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
         "Rancang Bangun Sistem Portal Berita\ndan Content Management System (CMS) Redaksi",
         font_size=32, bold=True, color=WHITE)
add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
         "di PT. Ketik Media Siber", font_size=20, color=GRAY)
add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
         "Oleh: Farrel Aqeel Danendra", font_size=18, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "NPM: 23081010204", font_size=14, color=GRAY)
add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.8),
         "Program Studi Informatika — Fakultas Ilmu Komputer\nUniversitas Pembangunan Nasional \"Veteran\" Jawa Timur — 2026",
         font_size=12, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: DAFTAR ISI
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(6), Inches(0.5),
         "DAFTAR ISI", font_size=28, bold=True, color=WHITE)

items = [
    ("01", "Gambaran Umum Proyek", "Profil perusahaan, tech stack, arsitektur sistem"),
    ("02", "Rekognisi: Analisis Kebutuhan", "Analisis kebutuhan fungsional & non-fungsional, Use Case, UI/UX"),
    ("03", "Rekognisi: Pemrograman API", "REST API endpoints, backend Laravel 12, frontend Next.js 15"),
    ("04", "Rekognisi: Uji Coba & Implementasi", "Black-box testing, pengujian fitur, dokumentasi hasil"),
    ("05", "Penutup", "Kesimpulan dan saran"),
]
y = Inches(1.5)
for num, title, desc in items:
    add_text(slide, Inches(1.0), y, Inches(0.8), Inches(0.5),
             num, font_size=24, bold=True, color=BLUE)
    add_text(slide, Inches(1.9), y, Inches(6), Inches(0.4),
             title, font_size=18, bold=True, color=DARK)
    add_text(slide, Inches(1.9), y + Inches(0.4), Inches(8), Inches(0.3),
             desc, font_size=12, color=GRAY)
    y += Inches(1.0)

slide_number_footer(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: GAMBARAN UMUM PROYEK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "GAMBARAN UMUM PROYEK", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.3), Inches(7), Inches(1.0),
         "Klojen.com adalah portal berita digital full-stack dengan arsitektur Headless CMS yang dikembangkan\n"
         "untuk PT. Ketik Media Siber, khusus melayani wilayah Kota Malang.",
         font_size=14, color=DARK)

# Tech stack cards
cards = [
    ("Backend", "Laravel 12\nPHP 8.2+\nMySQL 8\nJWT Auth"),
    ("Frontend", "Next.js 15\nTypeScript\nTailwind CSS\nZustand"),
    ("Arsitektur", "Headless CMS\nService-Repository\nREST API\nDecoupled"),
    ("Fitur Utama", "Editorial Workflow\nRBAC (4 Role)\nFull-Text Search\nSEO Optimized"),
]
x = Inches(0.8)
for title, content in cards:
    shape = add_rect(slide, x, Inches(3.0), Inches(2.7), Inches(3.5), LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), Inches(3.2), Inches(2.3), Inches(0.4),
             title, font_size=16, bold=True, color=BLUE)
    add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.3), Inches(2.5),
             content, font_size=12, color=DARK)
    x += Inches(3.0)

slide_number_footer(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: ARSITEKTUR SISTEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "ARSITEKTUR SISTEM", font_size=28, bold=True, color=WHITE)

# Architecture diagram using boxes
layers = [
    ("Frontend (Next.js 15)", Inches(1.3), BLUE, "Portal Publik + CMS Redaksi\nTailwind CSS, Zustand, Axios"),
    ("REST API (HTTP/JSON)", Inches(2.8), GREEN, "50 Endpoints\nBearer Token + JWT"),
    ("Backend (Laravel 12)", Inches(4.3), NAVY, "Controllers → Services → Repositories\nPHP 8.2+, Service-Repository Pattern"),
    ("Database (MySQL 8)", Inches(5.8), RGBColor(0x94,0x3B,0x00), "15 Tabel, UUID PK, FULLTEXT Index"),
]
for label, top, color, desc in layers:
    shape = add_rect(slide, Inches(1.5), top, Inches(10), Inches(1.1), color)
    add_text(slide, Inches(1.8), top + Inches(0.1), Inches(4), Inches(0.4),
             label, font_size=16, bold=True, color=WHITE)
    add_text(slide, Inches(6.5), top + Inches(0.1), Inches(4.5), Inches(0.8),
             desc, font_size=12, color=WHITE)

# Arrows between layers
for i in range(3):
    top = Inches(2.4) + Inches(i * 1.5)
    add_text(slide, Inches(6.2), top, Inches(1), Inches(0.4),
             "▼", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_image_placeholder(slide, Inches(0.5), Inches(1.2), Inches(0.8), Inches(6.0),
                      "[Diagram\nArsitektur]")
slide_number_footer(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: RBAC (Role-Based Access Control)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "ROLE-BASED ACCESS CONTROL (RBAC)", font_size=28, bold=True, color=WHITE)

roles = [
    ("Reader", "Membaca artikel, bookmark,\nkomentar, reset password", LIGHT_BLUE),
    ("Journalist", "Semua hak Reader +\nmenulis/edit artikel draft,\nupload media", RGBColor(0xDC,0xFC,0xE7)),
    ("Editor", "Semua hak Journalist +\nreview artikel, publish,\nlock/unlock, moderasi komentar", RGBColor(0xFE,0xF3,0xC7)),
    ("Admin", "Akses penuh: semua fitur\nCMS + kelola user, kategori,\ntag, dashboard analitik", RGBColor(0xFC,0xE7,0xE7)),
]
x = Inches(0.8)
for role, desc, bg in roles:
    shape = add_rect(slide, x, Inches(1.5), Inches(2.8), Inches(2.5), bg)
    add_text(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.5),
             role, font_size=20, bold=True, color=DARK)
    add_text(slide, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.5),
             desc, font_size=12, color=DARK)
    x += Inches(3.1)

# Editorial workflow
add_text(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.5),
         "Editorial Workflow (State Machine):", font_size=16, bold=True, color=DARK)

states = ["DRAFT", "REVIEW", "PUBLISHED", "SCHEDULED", "ARCHIVED"]
colors = [GRAY, ORANGE, GREEN, BLUE, RGBColor(0x94,0x3B,0x00)]
x = Inches(0.8)
for state, color in zip(states, colors):
    shape = add_rect(slide, x, Inches(5.2), Inches(1.8), Inches(0.8), color)
    add_text(slide, x + Inches(0.1), Inches(5.35), Inches(1.6), Inches(0.5),
             state, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if state != "ARCHIVED":
        add_text(slide, x + Inches(1.8), Inches(5.35), Inches(0.5), Inches(0.5),
                 "→", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
    x += Inches(2.3)

slide_number_footer(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: SECTION DIVIDER — REKOGNISI 1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "REKOGNISI 1", "Analisis Kebutuhan",
               "Identifikasi dan analisis kebutuhan sistem Portal Berita Klojen.com")
slide_number_footer(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: ANALISIS KEBUTUHAN FUNGSIONAL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "ANALISIS KEBUTUHAN FUNGSIONAL", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
         "Analisis dilakukan dengan mempelajari alur kerja redaksi PT. Ketik Media Siber secara menyeluruh.",
         font_size=14, color=GRAY)

needs = [
    ("Autentikasi & Otorisasi", "Register, login, logout, forgot/reset password\nRBAC: Reader, Journalist, Editor, Admin"),
    ("Portal Publik", "Beranda (featured, populer, terbaru), detail artikel,\nkomentar threaded, bookmark, pencarian FULLTEXT"),
    ("CMS Redaksi", "Dashboard berbasis role, Bank Berita, form tulis berita,\ntinjauan artikel, manajemen media, moderasi komentar"),
    ("Workflow Editorial", "State machine: draft → review → published/scheduled/\nrejected → archived, article locking, scheduled publishing"),
    ("Fitur Pendukung", "Manajemen media, full-text search, analitik page views,\nSEO (sitemap, JSON-LD, OG), rate limiting komentar"),
]
y = Inches(2.0)
for title, desc in needs:
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.9), BLUE)
    add_text(slide, Inches(1.1), y, Inches(3.5), Inches(0.4),
             title, font_size=14, bold=True, color=DARK)
    add_text(slide, Inches(1.1), y + Inches(0.35), Inches(6), Inches(0.5),
             desc, font_size=11, color=GRAY)
    y += Inches(1.05)

add_image_placeholder(slide, Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.5),
                      "[Screenshot halaman\nanalisis kebutuhan /\nFigma design brief]")
slide_number_footer(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: ANALISIS KEBUTUHAN NON-FUNGSIONAL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "ANALISIS KEBUTUHAN NON-FUNGSIONAL", font_size=28, bold=True, color=WHITE)

nf_needs = [
    ("Performa", "Optimasi SEO menggunakan SSG, lazy loading gambar, dan database indexing."),
    ("Keamanan", "JWT token pair + silent refresh, password hashing (bcrypt),\nrefresh token hashing (SHA-256), anti user-enumeration, article locking."),
    ("Skalabilitas", "Arsitektur Headless CMS — backend dan frontend\ndapat di-scale secara terpisah."),
    ("Responsivitas", "Antarmuka responsif untuk berbagai ukuran layar\nmenggunakan Tailwind CSS."),
    ("Pemeliharaan", "Service-Repository Pattern, dependency injection,\ndokumentasi teknis lengkap."),
]
y = Inches(1.5)
for title, desc in nf_needs:
    shape = add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.0), LIGHT_GRAY)
    add_text(slide, Inches(1.1), y + Inches(0.1), Inches(2.5), Inches(0.4),
             title, font_size=14, bold=True, color=BLUE)
    add_text(slide, Inches(3.8), y + Inches(0.1), Inches(8), Inches(0.8),
             desc, font_size=12, color=DARK)
    y += Inches(1.15)

slide_number_footer(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: USE CASE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "USE CASE DIAGRAM", font_size=28, bold=True, color=WHITE)

add_image_placeholder(slide, Inches(1.0), Inches(1.3), Inches(7.5), Inches(5.5),
                      "[MASUKKAN GAMBAR USE CASE DIAGRAM]")

add_text(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(0.4),
         "5 Level Aktor:", font_size=14, bold=True, color=DARK)
actors = [
    "Pengunjung — beranda, artikel, search",
    "Reader — komentar, bookmark, profil",
    "Journalist — tulis artikel, upload media",
    "Editor — review, publish, lock, moderasi",
    "Admin — kelola user, kategori, tag",
]
add_bullet_list(slide, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.0), actors, font_size=11)
slide_number_footer(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: PERANCANGAN UI/UX (FIGMA)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "PERANCANGAN UI/UX (FIGMA)", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Desain antarmuka dirancang menggunakan Figma sebelum diimplementasikan ke kode.",
         font_size=14, color=GRAY)

pages = [
    ("Halaman CMS", "Bank Berita, Tulis Berita,\nTinjauan Artikel, Preview Berita, Profil"),
    ("Portal Publik", "Beranda, Detail Artikel,\nHalaman Kategori"),
    ("Autentikasi", "Login, Register,\nReset Password"),
]
x = Inches(0.8)
for title, desc in pages:
    shape = add_rect(slide, x, Inches(1.9), Inches(3.8), Inches(1.4), LIGHT_BLUE)
    add_text(slide, x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4),
             title, font_size=14, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.7),
             desc, font_size=12, color=DARK)
    x += Inches(4.1)

# Image placeholders for Figma screenshots
add_image_placeholder(slide, Inches(0.8), Inches(3.6), Inches(3.8), Inches(3.3),
                      "[Screenshot Figma\nCMS Editor]")
add_image_placeholder(slide, Inches(4.9), Inches(3.6), Inches(3.8), Inches(3.3),
                      "[Screenshot Figma\nPortal Berita]")
add_image_placeholder(slide, Inches(9.0), Inches(3.6), Inches(3.8), Inches(3.3),
                      "[Screenshot Figma\nAutentikasi]")

slide_number_footer(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: SECTION DIVIDER — REKOGNISI 2
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "REKOGNISI 2", "Pemrograman API",
               "Pengembangan REST API backend (Laravel 12) dan integrasi frontend (Next.js 15)")
slide_number_footer(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: RINGKASAN API ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "RINGKASAN API ENDPOINTS", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Total 50 endpoint REST API yang dikelompokkan berdasarkan fungsi:",
         font_size=14, color=GRAY)

groups = [
    ("Auth", "9 endpoint", "register, login, refresh, logout, me,\nprofile, change-password, forgot/reset"),
    ("Portal Publik", "7 endpoint", "beranda, articles, detail, comments,\nkategori, tags, sitemap"),
    ("CMS Artikel", "8 endpoint", "CRUD artikel, update status,\nlock/unlock"),
    ("CMS Lainnya", "9 endpoint", "komentar moderasi, statistik,\nkategori, tag management"),
    ("Bookmark & Media", "6 endpoint", "toggle bookmark, upload media,\nlist media, delete media"),
    ("User Management", "6 endpoint", "CRUD user (admin only),\ndeactivate user"),
    ("Analytics", "1 endpoint", "track page view"),
]
x_start = Inches(0.5)
y = Inches(1.9)
for i, (group, count, desc) in enumerate(groups):
    col = i % 4
    row = i // 4
    x = x_start + Inches(col * 3.2)
    yy = y + Inches(row * 2.6)
    shape = add_rect(slide, x, yy, Inches(2.9), Inches(2.2), LIGHT_GRAY)
    add_text(slide, x + Inches(0.15), yy + Inches(0.1), Inches(2.6), Inches(0.3),
             group, font_size=13, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.15), yy + Inches(0.45), Inches(2.6), Inches(0.3),
             count, font_size=11, bold=True, color=BLUE)
    add_text(slide, x + Inches(0.15), yy + Inches(0.8), Inches(2.6), Inches(1.2),
             desc, font_size=10, color=DARK)

slide_number_footer(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: BACKEND — CREATE ARTICLE (CmsArticleService)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "BACKEND: CREATE ARTICLE", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
         "CmsArticleService::createArticle()", font_size=16, bold=True, color=DARK)
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.8),
         "Endpoint: POST /api/cms/articles\nTanggung jawab: Generate slug unik, INSERT articles (status=draft),\nbatch INSERT article_tags, INSERT initial revision, index ke search_indexes.",
         font_size=11, color=GRAY)

code = """public function createArticle(int $authorId, array $data): array
{
    // 1. Tentukan slug (auto-generate jika kosong)
    if (!empty($data['slug'])) {
        $slug = Str::slug($data['slug']);
        $this->assertSlugAvailable($slug);
    } else {
        $slug = $this->generateUniqueSlug($data['title']);
    }

    $articleId = (string) Str::uuid();

    // 2. INSERT articles (status = draft)
    DB::table('articles')->insert([
        'id'                 => $articleId,
        'author_id'          => $authorId,
        'category_id'        => $data['category_id'] ?? null,
        'title'              => $data['title'],
        'slug'               => $slug,
        'content'            => $data['content'],
        'featured_image_url' => $data['featured_image_url'] ?? null,
        'status'             => 'draft',
        'view_count'         => 0,
        'created_at'         => now(),
        'updated_at'         => now(),
    ]);

    // 3. Process tags + search index
    $this->processTags($articleId, $data['tag_ids'] ?? []);
    $this->searchService->reindexArticle($articleId);
}"""
add_code_block(slide, Inches(0.8), Inches(2.7), Inches(6.5), Inches(4.3), code,
               "backend/app/Services/CmsArticleService.php")

add_image_placeholder(slide, Inches(7.8), Inches(2.7), Inches(4.8), Inches(4.3),
                      "[Screenshot halaman\nTulis Berita (CMS)]")

slide_number_footer(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: BACKEND — EDIT PROFILE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "BACKEND: EDIT PROFILE", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
         "AuthController::updateProfile()", font_size=16, bold=True, color=DARK)
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
         "Endpoint: PUT /api/auth/profile\nValidasi: name (required), email (unique), avatar (nullable image, max 2MB)",
         font_size=11, color=GRAY)

code_profile = """public function updateProfile(Request $request): JsonResponse
{
    $user = auth('api')->user();

    $validated = $request->validate([
        'name'   => 'required|string|max:255',
        'email'  => 'required|string|email|max:255|unique:users,email,' . $user->id,
        'avatar' => 'nullable|image|mimes:jpeg,png,jpg|max:2048',
    ]);

    if ($request->hasFile('avatar')) {
        // Hapus avatar lama dari storage
        if ($user->avatar_url) {
            $oldPath = str_replace(asset('storage/'), '', $user->avatar_url);
            Storage::disk('public')->delete($oldPath);
        }
        $path = $request->file('avatar')->store('avatars', 'public');
        $validated['avatar_url'] = asset('storage/' . $path);
    }

    unset($validated['avatar']);
    $user->update($validated);

    return response()->json([
        'status'  => 'success',
        'message' => 'Profil berhasil diperbarui.',
    ]);
}"""
add_code_block(slide, Inches(0.8), Inches(2.4), Inches(6.5), Inches(4.6), code_profile,
               "backend/app/Http/Controllers/AuthController.php")

add_image_placeholder(slide, Inches(7.8), Inches(2.4), Inches(4.8), Inches(4.6),
                      "[Screenshot halaman\nEdit Profil]")

slide_number_footer(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: BACKEND — CHANGE PASSWORD
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "BACKEND: CHANGE PASSWORD", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
         "AuthController::changePassword()", font_size=16, bold=True, color=DARK)
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
         "Endpoint: PUT /api/auth/change-password\nValidasi: current_password (required), new_password (min 8 karakter)",
         font_size=11, color=GRAY)

code_pw = """public function changePassword(Request $request): JsonResponse
{
    $user = auth('api')->user();

    $validated = $request->validate([
        'current_password' => 'required|string',
        'new_password'     => 'required|string|min:8',
    ]);

    if (!Hash::check($validated['current_password'], $user->password)) {
        return response()->json([
            'status'  => 'error',
            'message' => 'Password saat ini salah.',
        ], 422);
    }

    $user->update([
        'password' => $validated['new_password'],
        // auto-hashed via Eloquent cast
    ]);

    return response()->json([
        'status'  => 'success',
        'message' => 'Password berhasil diubah.',
    ]);
}"""
add_code_block(slide, Inches(0.8), Inches(2.4), Inches(6.5), Inches(4.0), code_pw,
               "backend/app/Http/Controllers/AuthController.php")

add_image_placeholder(slide, Inches(7.8), Inches(2.4), Inches(4.8), Inches(4.0),
                      "[Screenshot halaman\nUbah Password]")

slide_number_footer(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: FRONTEND — AXIOS INTERCEPTOR & JWT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "FRONTEND: AXIOS INTERCEPTOR & JWT MANAGEMENT", font_size=28, bold=True, color=WHITE)

code_axios = """// frontend/lib/axios.ts
import axios from 'axios';
import { getAccessToken, setAccessToken } from '@/stores/authStore';
import { getRefreshToken, saveRefreshToken, clearRefreshToken } from '@/lib/auth';

const axiosInstance = axios.create({
  baseURL: process.env.NEXT_PUBLIC_API_URL,
  headers: { 'Content-Type': 'application/json' },
});

// Request interceptor — sisipkan JWT token otomatis
axiosInstance.interceptors.request.use((config) => {
  const token = getAccessToken();
  if (token) config.headers.Authorization = `Bearer ${token}`;
  return config;
});

// Response interceptor — silent refresh saat TOKEN_EXPIRED
axiosInstance.interceptors.response.use(
  (response) => response,
  async (error) => {
    const original = error.config;
    if (error.response?.data?.error === 'TOKEN_EXPIRED' && !original._retry) {
      original._retry = true;
      const res = await axiosInstance.post('/auth/refresh',
        { refresh_token: getRefreshToken() });
      setAccessToken(res.data.data.access_token);
      original.headers.Authorization = `Bearer ${res.data.data.access_token}`;
      return axiosInstance(original);
    }
    return Promise.reject(error);
  }
);"""
add_code_block(slide, Inches(0.8), Inches(1.3), Inches(7.0), Inches(5.8), code_axios,
               "frontend/lib/axios.ts")

add_text(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.4),
         "Mekanisme Silent Refresh:", font_size=14, bold=True, color=DARK)
steps = [
    "1. Request API → attach Bearer token",
    "2. Backend return TOKEN_EXPIRED",
    "3. Queue request yang gagal",
    "4. POST /auth/refresh + refresh_token",
    "5. Dapat access_token baru",
    "6. Retry request awal dengan token baru",
    "7. Jika refresh gagal → force logout",
]
add_bullet_list(slide, Inches(8.3), Inches(2.0), Inches(4.5), Inches(3.5), steps, font_size=11)

add_image_placeholder(slide, Inches(8.3), Inches(5.0), Inches(4.5), Inches(2.0),
                      "[Screenshot Axios config\ndi VS Code / browser]")

slide_number_footer(slide, 16, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: FRONTEND — KOMPONEN REUSABLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "FRONTEND: KOMPONEN REUSABLE", font_size=28, bold=True, color=WHITE)

components = [
    ("Rich Text Editor", "Editor konten artikel menggunakan\nlibrary React (TipTap). Mendukung\nbold, italic, heading, list, link,\ndan penyisipan gambar."),
    ("Image Cropper", "Komponen crop gambar menggunakan\nreact-easy-crop. Digunakan untuk\navatar profil dan featured image.\nHasil berupa Blob → FormData."),
    ("Skeleton Loading", "Placeholder animasi saat data\nloading. Menggunakan animate-pulse\nTailwind CSS. Setiap halaman\npunya skeleton custom."),
    ("Toast Notification", "Popup feedback tanpa library\neksternal. Tipe: success, error,\nwarning. Auto-dismiss 4 detik."),
    ("Modal Konfirmasi", "Dialog popup untuk aksi destruktif.\nBackdrop blur, fade-in, tombol\nkonfirmasi & batal."),
    ("Watermark Canvas", "Bake watermark \"KLOJEN.COM\"\nke gambar via Canvas API sebelum\nupload. Rotasi -30°, opacity 40%."),
]
x = Inches(0.5)
y = Inches(1.3)
for i, (title, desc) in enumerate(components):
    col = i % 3
    row = i // 3
    cx = x + Inches(col * 4.2)
    cy = y + Inches(row * 2.9)
    shape = add_rect(slide, cx, cy, Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(3.4), Inches(0.4),
             title, font_size=14, bold=True, color=NAVY)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.6), Inches(3.4), Inches(1.7),
             desc, font_size=11, color=DARK)

slide_number_footer(slide, 17, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: FRONTEND — HALAMAN TULIS BERITA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "FRONTEND: HALAMAN TULIS BERITA", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.8),
         "Fitur Utama:\n- Rich Text Editor (TipTap)\n- Upload featured image + crop\n- Watermark otomatis (Canvas API)\n- Pilih kategori & tag\n- Multi-foto (max 3)\n- Submit sebagai Draft / Review / Publish",
         font_size=12, color=DARK)

add_image_placeholder(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.5),
                      "[Screenshot halaman\nTulis Berita]")
add_image_placeholder(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.5),
                      "[Screenshot Rich Text\nEditor / Image Cropper]")

slide_number_footer(slide, 18, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: FRONTEND — EDIT PROFIL & CHANGE PASSWORD
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "FRONTEND: EDIT PROFIL & UBAH PASSWORD", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
         "Edit Profil (PUT /api/auth/profile)", font_size=14, bold=True, color=BLUE)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(1.0),
         "- Edit nama, email, avatar\n- Upload avatar + crop (react-easy-crop)\n- FormData + _method=PUT trick untuk Laravel\n- Fallback fetch /auth/me jika state kosong",
         font_size=11, color=DARK)

add_text(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(0.4),
         "Ubah Password (PUT /api/auth/change-password)", font_size=14, bold=True, color=BLUE)
add_text(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(1.0),
         "- Validasi: password lama, baru, konfirmasi\n- Minimal 8 karakter\n- Toggle show/hide password\n- Error handling dari backend",
         font_size=11, color=DARK)

add_image_placeholder(slide, Inches(6.5), Inches(1.2), Inches(6.0), Inches(2.8),
                      "[Screenshot halaman\nEdit Profil]")
add_image_placeholder(slide, Inches(6.5), Inches(4.3), Inches(6.0), Inches(2.8),
                      "[Screenshot halaman\nUbah Password]")

slide_number_footer(slide, 19, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: SECTION DIVIDER — REKOGNISI 3
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "REKOGNISI 3", "Uji Coba dan Implementasi",
               "Pengujian fitur menggunakan metode Black-Box Testing")
slide_number_footer(slide, 20, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21: METODE PENGUJIAN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "METODE PENGUJIAN: BLACK-BOX TESTING", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Pengujian dilakukan menggunakan metode Black-Box Testing — menguji fungsionalitas sistem tanpa\n"
         "melihat struktur kode internal. Setiap fitur diuji dengan skenario input valid, invalid, dan edge case.",
         font_size=14, color=DARK)

test_groups = [
    ("Autentikasi", "Register, Login, Logout, Refresh Token,\nForgot/Reset Password"),
    ("Portal Publik", "Beranda, Detail Artikel, Search,\nKategori, Bookmark, Komentar"),
    ("CMS Redaksi", "CRUD Artikel, Update Status, Lock/Unlock,\nTinjauan Artikel, Media, Moderasi Komentar"),
    ("Editorial Workflow", "Draft → Review → Published/Scheduled/Rejected\n→ Archived, Scheduled Publish (cron)"),
    ("Profil & Keamanan", "Edit Profil, Ubah Password, Upload Avatar,\nJWT Silent Refresh, Role Guard"),
]
x = Inches(0.5)
y = Inches(2.5)
for i, (title, desc) in enumerate(test_groups):
    col = i % 3
    row = i // 3
    cx = x + Inches(col * 4.2)
    cy = y + Inches(row * 2.3)
    shape = add_rect(slide, cx, cy, Inches(3.8), Inches(1.8), LIGHT_BLUE)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(3.4), Inches(0.4),
             title, font_size=14, bold=True, color=NAVY)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.6), Inches(3.4), Inches(1.0),
             desc, font_size=11, color=DARK)

slide_number_footer(slide, 21, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22: HASIL PENGUJIAN — AUTENTIKASI & PORTAL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "HASIL PENGUJIAN: AUTENTIKASI & PORTAL PUBLIK", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
         "Testing Autentikasi", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.0),
                      "[Tabel testing autentikasi:\nlogin, register, refresh, logout]")

add_text(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
         "Testing Portal Publik", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
                      "[Tabel testing portal:\nberanda, artikel, search, bookmark, komentar]")

add_text(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4),
         "Screenshot Pengujian", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.3),
                      "[Screenshot Postman / browser\nhasil pengujian autentikasi\ndan portal publik]")

slide_number_footer(slide, 22, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23: HASIL PENGUJIAN — CMS & EDITORIAL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
add_text(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
         "HASIL PENGUJIAN: CMS & EDITORIAL WORKFLOW", font_size=28, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
         "Testing CMS Redaksi", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.0),
                      "[Tabel testing CMS:\nCRUD artikel, media, moderasi komentar]")

add_text(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
         "Testing Editorial Workflow", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
                      "[Tabel testing workflow:\ndraft→review→publish/schedule/reject→archive]")

add_text(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4),
         "Screenshot Pengujian", font_size=16, bold=True, color=BLUE)
add_image_placeholder(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.3),
                      "[Screenshot Postman / browser\nhasil pengujian CMS\ndan editorial workflow]")

slide_number_footer(slide, 23, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24: KESIMPULAN & PENUTUP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
         "KESIMPULAN", font_size=32, bold=True, color=WHITE)

conclusions = [
    "Portal Berita Klojen.com berhasil dibangun menggunakan arsitektur Headless CMS\ndengan Laravel 12 (backend) dan Next.js 15 (frontend).",
    "Rekognisi Analisis Kebutuhan — Analisis kebutuhan fungsional dan non-fungsional\nberhasil dilakukan melalui studi alur kerja redaksi PT. Ketik Media Siber.",
    "Rekognisi Pemrograman API — 50 endpoint REST API berhasil dikembangkan,\ntermasuk fitur create article, edit profile, dan change password.",
    "Rekognisi Uji Coba & Implementasi — Seluruh fitur berhasil diuji menggunakan\nmetode Black-Box Testing dan berfungsi sesuai spesifikasi.",
    "Sistem siap digunakan untuk mendukung operasional redaksi Klojen.com\ndengan workflow editorial yang terstruktur dan aman.",
]
y = Inches(2.0)
for item in conclusions:
    add_rect(slide, Inches(1.5), y, Inches(0.06), Inches(0.6), BLUE)
    add_text(slide, Inches(1.8), y - Inches(0.05), Inches(9), Inches(0.7),
             item, font_size=13, color=WHITE)
    y += Inches(1.0)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "TERIMA KASIH", font_size=24, bold=True, color=BLUE)
slide_number_footer(slide, 24, TOTAL_SLIDES)

# ── Save ─────────────────────────────────────────────────────────────────────
output_path = r"c:\apps\klojen.com\Presentasi_PLK2_Farrel.pptx"
prs.save(output_path)
print(f"✅ PPT berhasil dibuat: {output_path}")
print(f"   Total slide: {len(prs.slides)}")
